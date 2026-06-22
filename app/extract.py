"""
① 입력 정규화 모듈
PDF / PPTX 를 "페이지 = 구조화 텍스트 (+ PDF는 페이지 이미지)" 형태로 통일한다.

- PDF : pdfplumber(텍스트, layout 보존) + PyMuPDF(fitz, 페이지 PNG 렌더)
- PPTX: python-pptx (표 셀 / 텍스트박스 구조 그대로). 이미지 렌더는 하지 않음(가벼운 구성).

반환:
  pages: List[Page]  (1-base page_no, text, image_path|None, kind)
외부 바이너리(pdftotext/libreoffice) 의존 없음.
"""
from __future__ import annotations

import os
from dataclasses import dataclass, field
from typing import Optional


@dataclass
class Page:
    page_no: int                     # 1-base
    text: str
    image_path: Optional[str] = None  # PDF만 채워짐 (검토 게이트 대조용)
    kind: str = "pdf"                 # "pdf" | "pptx"


@dataclass
class Document:
    pages: list[Page] = field(default_factory=list)
    kind: str = "pdf"

    @property
    def page_count(self) -> int:
        return len(self.pages)

    def text_index(self, max_chars_per_page: int = 1800) -> str:
        """②인덱싱용: 페이지 번호 + 앞부분 텍스트만 모은 경량 인덱스."""
        out = []
        for p in self.pages:
            snippet = (p.text or "").strip().replace("\u0000", "")
            if len(snippet) > max_chars_per_page:
                snippet = snippet[:max_chars_per_page] + " …(생략)"
            out.append(f"[PAGE {p.page_no}]\n{snippet}")
        return "\n\n".join(out)


class CorruptDocumentError(Exception):
    """PDF/PPTX 파일이 손상/잘림 등으로 열리지 않을 때."""


def extract_pdf(path: str, image_dir: str, dpi: int = 110) -> Document:
    import pdfplumber
    import fitz  # PyMuPDF

    doc = Document(kind="pdf")
    os.makedirs(image_dir, exist_ok=True)

    size = os.path.getsize(path) if os.path.exists(path) else 0

    # 텍스트
    texts: dict[int, str] = {}
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                t = page.extract_text(layout=True) or page.extract_text() or ""
                texts[i] = t
    except Exception as e:  # PdfminerException(Unexpected EOF) 등
        raise CorruptDocumentError(
            f"PDF를 열 수 없습니다(파일 손상/업로드 잘림 가능). "
            f"현재 크기 {size:,}바이트. 원본 파일을 다시 업로드해 주세요. (원인: {e})"
        ) from e

    # 페이지 이미지 (PNG)
    try:
        fdoc = fitz.open(path)
    except Exception as e:
        raise CorruptDocumentError(
            f"PDF 페이지 렌더링에 실패했습니다(파일 손상/업로드 잘림 가능). "
            f"현재 크기 {size:,}바이트. 원본 파일을 다시 업로드해 주세요. (원인: {e})"
        ) from e
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    for i in range(fdoc.page_count):
        page_no = i + 1
        pix = fdoc.load_page(i).get_pixmap(matrix=mat)
        img_path = os.path.join(image_dir, f"page_{page_no:04d}.png")
        pix.save(img_path)
        doc.pages.append(Page(
            page_no=page_no,
            text=texts.get(page_no, ""),
            image_path=img_path,
            kind="pdf",
        ))
    fdoc.close()
    return doc


def extract_pptx(path: str) -> Document:
    from pptx import Presentation

    doc = Document(kind="pptx")
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(" | ".join(cells))
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip() or para.text.strip()
                    if t:
                        lines.append(t)
        doc.pages.append(Page(
            page_no=i,
            text="\n".join(lines),
            image_path=None,   # 가벼운 구성: PPTX 이미지 렌더 없음
            kind="pptx",
        ))
    return doc


def normalize(path: str, image_dir: str) -> Document:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path, image_dir)
    if ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    raise ValueError(f"지원하지 않는 형식: {ext} (PDF 또는 PPTX만 가능)")
